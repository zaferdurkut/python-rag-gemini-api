import io
import logging
from typing import Dict, Any, Optional, List
from pathlib import Path
import mimetypes

# File processing imports with optional dependencies
try:
    import pypdf

    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl

    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract

    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

logger = logging.getLogger(__name__)


class FileProcessor:
    """Service for processing various file types and extracting text content."""

    SUPPORTED_EXTENSIONS = {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".doc": "application/msword",
        ".txt": "text/plain",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xls": "application/vnd.ms-excel",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".ppt": "application/vnd.ms-powerpoint",
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".bmp": "image/bmp",
        ".tiff": "image/tiff",
    }

    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    MAX_IMAGE_SIZE = 10 * 1024 * 1024  # 10MB for images

    def __init__(self):
        """Initialize the file processor."""
        self.supported_types = list(self.SUPPORTED_EXTENSIONS.keys())

    def is_supported_file(self, filename: str) -> bool:
        """Check if file type is supported."""
        file_ext = Path(filename).suffix.lower()
        return file_ext in self.SUPPORTED_EXTENSIONS

    def get_file_type(self, filename: str) -> str:
        """Get the file type from filename."""
        file_ext = Path(filename).suffix.lower()
        return self.SUPPORTED_EXTENSIONS.get(file_ext, "application/octet-stream")

    def validate_file_size(self, file_size: int, is_image: bool = False) -> bool:
        """Validate file size."""
        max_size = self.MAX_IMAGE_SIZE if is_image else self.MAX_FILE_SIZE
        return file_size <= max_size

    async def process_file(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process a file and extract text content."""
        try:
            file_ext = Path(filename).suffix.lower()

            if not self.is_supported_file(filename):
                raise ValueError(f"Unsupported file type: {file_ext}")

            # Check file size
            file_size = len(file_content)
            is_image = file_ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"]
            if not self.validate_file_size(file_size, is_image):
                max_size = self.MAX_IMAGE_SIZE if is_image else self.MAX_FILE_SIZE
                raise ValueError(
                    f"File too large. Max size: {max_size / (1024*1024):.1f}MB"
                )

            # Extract text based on file type
            text_content = await self._extract_text(file_content, file_ext)

            # Create metadata
            metadata = {
                "filename": filename,
                "file_type": self.get_file_type(filename),
                "file_size": file_size,
                "processed_at": self._get_timestamp(),
                "text_length": len(text_content),
            }

            return {"content": text_content, "metadata": metadata, "success": True}

        except Exception as e:
            logger.error(f"Error processing file {filename}: {e}")
            return {
                "content": "",
                "metadata": {
                    "filename": filename,
                    "error": str(e),
                    "processed_at": self._get_timestamp(),
                },
                "success": False,
            }

    async def _extract_text(self, file_content: bytes, file_ext: str) -> str:
        """Extract text content from file based on type."""
        try:
            if file_ext == ".pdf":
                return await self._extract_pdf_text(file_content)
            elif file_ext in [".docx"]:
                return await self._extract_docx_text(file_content)
            elif file_ext in [".txt"]:
                return file_content.decode("utf-8")
            elif file_ext in [".xlsx", ".xls"]:
                return await self._extract_excel_text(file_content)
            elif file_ext in [".pptx", ".ppt"]:
                return await self._extract_powerpoint_text(file_content)
            elif file_ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"]:
                return await self._extract_image_text(file_content)
            else:
                raise ValueError(f"Unsupported file type: {file_ext}")

        except Exception as e:
            logger.error(f"Error extracting text from {file_ext}: {e}")
            raise

    async def _extract_pdf_text(self, file_content: bytes) -> str:
        """Extract text from PDF file."""
        if not PDF_AVAILABLE:
            raise ValueError("PDF processing not available. Install pypdf package.")

        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = pypdf.PdfReader(pdf_file)

            text_parts = []
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text.strip():
                    text_parts.append(text.strip())

            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            raise

    async def _extract_docx_text(self, file_content: bytes) -> str:
        """Extract text from DOCX file."""
        if not DOCX_AVAILABLE:
            raise ValueError(
                "DOCX processing not available. Install python-docx package."
            )

        try:
            doc_file = io.BytesIO(file_content)
            doc = Document(doc_file)

            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text.strip())

            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            raise

    async def _extract_excel_text(self, file_content: bytes) -> str:
        """Extract text from Excel file."""
        if not EXCEL_AVAILABLE:
            raise ValueError(
                "Excel processing not available. Install openpyxl package."
            )

        try:
            excel_file = io.BytesIO(file_content)
            workbook = openpyxl.load_workbook(excel_file)

            text_parts = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = []

                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        sheet_text.append(" | ".join(row_text))

                if sheet_text:
                    text_parts.append(f"Sheet: {sheet_name}\n" + "\n".join(sheet_text))

            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting Excel text: {e}")
            raise

    async def _extract_powerpoint_text(self, file_content: bytes) -> str:
        """Extract text from PowerPoint file."""
        if not PPTX_AVAILABLE:
            raise ValueError(
                "PowerPoint processing not available. Install python-pptx package."
            )

        try:
            ppt_file = io.BytesIO(file_content)
            presentation = Presentation(ppt_file)

            text_parts = []
            for i, slide in enumerate(presentation.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if slide_text:
                    text_parts.append(f"Slide {i}:\n" + "\n".join(slide_text))

            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting PowerPoint text: {e}")
            raise

    async def _extract_image_text(self, file_content: bytes) -> str:
        """Extract text from image using OCR."""
        if not OCR_AVAILABLE:
            raise ValueError(
                "OCR processing not available. Install Pillow and pytesseract packages."
            )

        try:
            image = Image.open(io.BytesIO(file_content))

            # Convert to RGB if necessary
            if image.mode != "RGB":
                image = image.convert("RGB")

            # Use pytesseract for OCR
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting image text: {e}")
            raise

    def _get_timestamp(self) -> str:
        """Get current timestamp."""
        from datetime import datetime

        return datetime.now().isoformat()

    async def process_multiple_files(
        self, files: List[Dict[str, Any]]
    ) -> List[Dict[str, Any]]:
        """Process multiple files."""
        results = []
        for file_data in files:
            result = await self.process_file(
                file_data["content"], file_data["filename"]
            )
            results.append(result)
        return results
